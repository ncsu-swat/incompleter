# LExecutor: DO NOT INSTRUMENT

#Source: https://stackoverflow.com/questions/70286125/how-to-fix-attributeerror-image-object-has-no-attribute-seek-did-you-mean
from lexecutor.Runtime import _n_
from lexecutor.Runtime import _a_
from lexecutor.Runtime import _c_
from lexecutor.Runtime import _l_
try:
    from io import FileIO
    _l_(555983)

except ImportError:
    pass
try:
    import os
    _l_(555985)

except ImportError:
    pass
try:
    from wand.image import Image
    _l_(555987)

except ImportError:
    pass
try:
    from pptx.util import Inches
    _l_(555989)

except ImportError:
    pass
try:
    from pptx import Presentation
    _l_(555991)

except ImportError:
    pass

def create_slide()->_n_(555992, "FileIO", lambda: FileIO):
    _l_(556069)

    # Creating presentation object
    root = _c_(555994, _n_(555993, "Presentation", lambda: Presentation))
    _l_(555995)
    for file in _c_(555998, _a_(555997, _n_(555996, "os", lambda: os), "listdir")):
        _l_(556068)

        if _c_(556001, _a_(556000, _n_(555999, "file", lambda: file), "startswith"), 'watermarked_'):
            _l_(556067)

            # Creating slide layout
            first_slide_layout = _a_(556003, _n_(556002, "root", lambda: root), "slide_layouts")[1] 
            _l_(556004) 
            slide = _c_(556009, _a_(556007, _a_(556006, _n_(556005, "root", lambda: root), "slides"), "add_slide"), _n_(556008, "first_slide_layout", lambda: first_slide_layout))
            _l_(556010)
            shapes = _a_(556012, _n_(556011, "slide", lambda: slide), "shapes")
            _l_(556013)
            
            #Adding title or heading to the slide
            title_shape = _a_(556015, _n_(556014, "shapes", lambda: shapes), "title")
            _l_(556016)
            _n_(556017, "title_shape", lambda: title_shape).text = f" Created By python-pptx for Watermarking "
            _l_(556018)
            
            #Adding sub-title with border to the slide
            body_shape = _a_(556020, _n_(556019, "shapes", lambda: shapes), "placeholders")[1]
            _l_(556021)
            tf = _a_(556023, _n_(556022, "body_shape", lambda: body_shape), "text_frame")
            _l_(556024)
            _n_(556025, "tf", lambda: tf).text = f"This is a watermarked image of {_n_(556026, 'file', lambda: file)}"
            _l_(556027)
                
            with _c_(556030, _n_(556028, "Image", lambda: Image), filename = _n_(556029, "file", lambda: file)) as watermarked_image:
                _l_(556066)

                
                #Maintianing the aspect ratio of the image
                width, height = _a_(556032, _n_(556031, "watermarked_image", lambda: watermarked_image), "size")
                _l_(556033)
                ratio = _n_(556034, "height", lambda: height)/_n_(556035, "width", lambda: width)
                _l_(556036)
                new_width = _n_(556037, "width", lambda: width) / 2
                _l_(556038)
                new_height = _c_(556042, _n_(556039, "int", lambda: int), _n_(556040, "new_width", lambda: new_width) * _n_(556041, "ratio", lambda: ratio))
                _l_(556043)
                _c_(556050, _a_(556045, _n_(556044, "watermarked_image", lambda: watermarked_image), "resize"), _c_(556048, _n_(556046, "int", lambda: int), _n_(556047, "new_width", lambda: new_width)), _n_(556049, "new_height", lambda: new_height))
                _l_(556051)
                
                # Add the watermarked image to the slide
                _c_(556060, _a_(556054, _a_(556053, _n_(556052, "slide", lambda: slide), "shapes"), "add_picture"), _n_(556055, "watermarked_image", lambda: watermarked_image) ,_c_(556057, _n_(556056, "Inches", lambda: Inches), 1), _c_(556059, _n_(556058, "Inches", lambda: Inches), 3))
                _l_(556061)
                _c_(556064, _a_(556063, _n_(556062, "root", lambda: root), "save"), "Output.pptx")
                _l_(556065)

_c_(556071, _n_(556070, "create_slide", lambda: create_slide))
_l_(556072)