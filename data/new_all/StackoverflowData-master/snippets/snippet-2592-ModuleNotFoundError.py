#Source: https://stackoverflow.com/questions/70286125/how-to-fix-attributeerror-image-object-has-no-attribute-seek-did-you-mean
from io import FileIO
import os
from wand.image import Image
from pptx.util import Inches 
from pptx import Presentation 

def create_slide()->FileIO:
    # Creating presentation object
    root = Presentation()
    for file in os.listdir():
        if file.startswith('watermarked_'):
            # Creating slide layout
            first_slide_layout = root.slide_layouts[1] 
            slide = root.slides.add_slide(first_slide_layout)
            shapes = slide.shapes
            
            #Adding title or heading to the slide
            title_shape = shapes.title
            title_shape.text = f" Created By python-pptx for Watermarking "
            
            #Adding sub-title with border to the slide
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.text = f"This is a watermarked image of {file}"
                
            with Image(filename = file) as watermarked_image:
                
                #Maintianing the aspect ratio of the image
                width, height = watermarked_image.size
                ratio = height/width
                new_width = width / 2
                new_height = int(new_width * ratio)
                watermarked_image.resize(int(new_width), new_height)
                
                # Add the watermarked image to the slide
                slide.shapes.add_picture(watermarked_image ,Inches(1), Inches(3))
                root.save("Output.pptx")

create_slide()